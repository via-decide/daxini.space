# make_daxini_pptx.py
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor

prs = Presentation()
# Basic master settings
PRIMARY = RGBColor(0x1F, 0x2A, 0x44)   # #1F2A44
ACCENT = RGBColor(0xFF, 0x8A, 0x00)    # #FF8A00
NEUTRAL = RGBColor(0xF5, 0xF7, 0xFA)   # #F5F7FA

slides = [
    {
        "title": "The Daxini Stack — Free 4‑Day Download Window",
        "bullets": [
            "All five volumes of sovereign AI architecture — EPUB, free starting today",
            "dharam · Daxini.space/books · #SovereignAI #IndiaTech"
        ],
        "notes": "Announce the limited window and the core promise: free access to five technical volumes and entry into an open ecosystem. Emphasize no email required and immediate download.",
    },
    {
        "title": "What's Free Right Now",
        "bullets": [
            "All five volumes of the Daxini Stack in EPUB format.",
            "Free to download for 4 days starting today.",
            "No email required; immediate access at Daxini.space/books."
        ],
        "notes": "Keep this short and urgent — the price is zero, the scarcity is early access and community membership. Repeat the URL slowly.",
    },
    {
        "title": "Not Just Books — An Ecosystem Entry Point",
        "bullets": [
            "Books act as on‑ramps into tools, communities, and early product access.",
            "Download → Read → Choose your path: Builder, Founder, Student.",
            "Volume 6 (beta chapters) available to early community members only."
        ],
        "notes": "Stress that the limited window is about joining early, not a discount. Early members shape Volume 6 and the tooling.",
    },
    {
        "title": "Choose Your Path After Reading",
        "bullets": [
            "Builder: Zayvora reasoning engine; LogicHub decision framework; vendor‑free orchestration tools.",
            "Founder: SOP systems; distribution frameworks; scaling patterns from 0 → 10,000+ readers.",
            "Student: Interactive learning systems; exam prep frameworks; architecture walkthroughs."
        ],
        "notes": "Explain each path briefly and how the books map to practical next steps and tools. Mention that Volume 6 deepens technical content for builders.",
    },
    {
        "title": "For Builders — Tools and Architecture",
        "bullets": [
            "Zayvora: reasoning engine for local and sovereign inference.",
            "LogicHub: decision framework for deterministic, auditable flows.",
            "Orchestration tools: run agentic systems without cloud vendor lock‑in."
        ],
        "notes": "Highlight technical depth in Volumes 1–5 and how Volume 6 will include deep dives and beta access to Zayvora APIs. Mention examples of use cases (edge inference, private LLM orchestration).",
    },
    {
        "title": "For Founders — Systems to Scale",
        "bullets": [
            "SOP systems for reproducible operations.",
            "Distribution frameworks for organic growth and community‑driven reach.",
            "Scaling patterns: reader acquisition, retention, and monetization playbooks."
        ],
        "notes": "Emphasize playbooks in the books and community case studies that show 0→10k reader growth patterns. Call out practical templates founders can reuse.",
    },
    {
        "title": "For Students — Learn by Building",
        "bullets": [
            "Interactive architecture walkthroughs and labs.",
            "Exam prep frameworks and practice problems.",
            "Mentorship and study groups inside the ecosystem."
        ],
        "notes": "Explain how students can use the EPUBs plus community resources to build portfolios and prepare for interviews. Mention opportunities for project-based learning.",
    },
    {
        "title": "Volume 6 Early Access Mechanics",
        "bullets": [
            "Download Volumes 1–5 now. No signup required.",
            "Opt in inside the ecosystem for Volume 6 beta chapters.",
            "Early access closes when the 4‑day window ends — community membership remains free."
        ],
        "notes": "Clarify that early access is time‑limited and tied to community participation and feedback. Explain how to opt in after downloading.",
    },
    {
        "title": "Join 10,000+ Engineers and Founders",
        "bullets": [
            "Active builders, founders, and students collaborating on sovereign AI.",
            "Share the books, use the tools, give feedback — that's how the ecosystem grows.",
            "Hashtags and channels: #SovereignAI #IndiaTech #BuildInPublic #Daxini."
        ],
        "notes": "Use this slide to invite social sharing and to show that momentum already exists. Mention any notable contributors or early adopters if applicable.",
    },
    {
        "title": "How to Get the Books and Join",
        "bullets": [
            "Visit: Daxini.space/books — download all five EPUBs now.",
            "No email required; immediate download.",
            "Share, use the tools, and opt into Volume 6 beta if you want early technical access."
        ],
        "notes": "Repeat the URL slowly and clearly. Encourage attendees to download now while the window is open and to invite peers.",
    },
    {
        "title": "Timeline and What Happens Next",
        "bullets": [
            "Day 0–4: Free download window open.",
            "After window: Volume 6 beta closes to new early access members.",
            "Ongoing: Community events, tool betas, and public roadmaps."
        ],
        "notes": "Reinforce urgency and the benefit of joining early to influence Volume 6 and tooling. Mention upcoming community calls or office hours if scheduled.",
    },
    {
        "title": "Thank You and Stay Connected",
        "bullets": [
            "Download: Daxini.space/books",
            "Follow: #SovereignAI #IndiaTech #BuildInPublic #Daxini",
            "Contact: community@daxini.space or site contact form"
        ],
        "notes": "Close with a short invitation to ask questions and to join the community channels. Offer to share the slide deck and one‑page summary after the talk.",
    },
]

for s in slides:
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = s["title"]
    tf = slide.shapes.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(s["bullets"]):
        if i == 0:
            p = tf.paragraphs[0]
            p.text = b
            p.font.size = Pt(20)
        else:
            p = tf.add_paragraph()
            p.text = b
            p.level = 1
            p.font.size = Pt(18)
    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = s["notes"]

# Save file
outname = "Daxini_Stack_Free_4Day_Window.pptx"
prs.save(outname)
print("Saved:", outname)
